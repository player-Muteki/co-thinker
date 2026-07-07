from __future__ import annotations

import io
from abc import ABC, abstractmethod
from pathlib import Path


class DocumentParser(ABC):
    @abstractmethod
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        ...


class TextParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        encodings = ("utf-8", "utf-8-sig", "gb18030", "latin-1")
        for encoding in encodings:
            try:
                return file_bytes.decode(encoding)
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Unable to decode file with any of {encodings}")


class PDFParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        import fitz

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages: list[str] = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text.strip():
                pages.append(text)
        doc.close()
        return "\n\n".join(pages)


class DocxParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        paragraphs: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        return "\n\n".join(paragraphs)


class PPTXParser(DocumentParser):
    def parse(self, file_bytes: bytes, file_ext: str) -> str:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        slides: list[str] = []
        for slide in prs.slides:
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(row_cells))
            if slide_texts:
                slides.append("\n".join(slide_texts))
        return "\n\n".join(slides)


TEXT_EXTENSIONS = frozenset({
    ".c", ".cpp", ".cs", ".go", ".h", ".java", ".js", ".jsx",
    ".md", ".mdx", ".php", ".py", ".rb", ".rs", ".ts", ".tsx", ".txt",
})


def _build_default_registry() -> dict[str, DocumentParser]:
    text_parser = TextParser()
    pdf_parser = PDFParser()
    docx_parser = DocxParser()
    pptx_parser = PPTXParser()

    registry: dict[str, DocumentParser] = {}
    for ext in TEXT_EXTENSIONS:
        registry[ext] = text_parser
    registry[".pdf"] = pdf_parser
    registry[".docx"] = docx_parser
    registry[".pptx"] = pptx_parser
    return registry


PARSER_REGISTRY: dict[str, DocumentParser] = _build_default_registry()
